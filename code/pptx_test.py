from pptx import Presentation

path = '../extras/GUI Test.pptx'
# load a presentation
prs = Presentation(path)
slide = prs.slides[0]
shapes = slide.shapes
print(shapes)